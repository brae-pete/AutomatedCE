"""
This file contains all the functions/methods to write data from the System to a file object.

It also contains all functions to read in data from the System
"""
from L3.SystemsBuilder import CESystem
from L4 import Electropherogram

from scipy.interpolate import interp1d
import os
from pathlib import Path
import numpy as np
from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches
from pptx.util import Pt
from scipy.signal import find_peaks
import pandas as pd

# Get the Home Directory Location
HOME = os.getcwd().split("AutomatedCE")[0] + "AutomatedCE"

def get_system_var(*var_names):
    """
    Get a variable from the system config file
    :param var_names: str
    :return:
    """
    # Get the Var file
    var_path = os.path.join(HOME, "config", "system_var.txt")
    with open(var_path) as fin:
        var_lines = fin.readlines()

    var_dict = {}

    for var_str in var_lines:
        var_list = var_str.split(',')
        var_dict[var_list[0]] = eval(var_list[1].replace('\n', ''))

    response = []
    for var_name in var_names:
        assert var_name in var_dict.keys(), f"{var_name} is not a var in system_var.txt"
        response.append(var_dict[var_name])

    return response


def get_data_filename(name, file_dir, extension=".csv"):
    """
    Returns a unique filename by adding 0's to the end of the name (but before the file suffix).
    :param name: str
    :param file_dir: str
    :return file_name: str
    """
    # Check if directory exists and make it if not
    Path(file_dir).mkdir(parents=True, exist_ok=True)

    file_name = os.path.join(file_dir, name + extension)
    copy_number = 0
    while os.path.exists(file_name):
        new_name = name + f"_{copy_number:05d}" + extension
        file_name = os.path.join(file_dir, new_name)
        copy_number += 1
    return file_name


def get_data_folder(name, file_dir):
    """
    Returns a unique folder name by adding 0's to the end of the name
    :param name:
    :param file_dir:
    :return:
    """
    # Check if directory exists and make it if not
    Path(file_dir).mkdir(parents=True, exist_ok=True)

    folder = os.path.join(file_dir, name)
    copy_number = 0
    while os.path.isdir(folder):
        new_name = name + f"_{copy_number:05d}"
        folder = os.path.join(file_dir, new_name)
        copy_number += 1
    Path(folder).mkdir(parents=True, exist_ok=True)
    return folder


class OutputElectropherogram:
    """
    Writes electropherogram data from the system to the
    """

    def __init__(self, system: CESystem, filepath, header=""):
        """
        Get the data from the System and write it to the file specified
        :param system:
        :param filepath:
        """
        data = self.get_data(system)
        self.write_data(data, header, filepath)

    @staticmethod
    def write_data(data, header, filepath):
        """
        Writes the CE data file containing a time_data column, a RFU, a kV and a uA column.
        incoming data should be in that order (time_data, rfu, voltage, current)
        data = [list of lists or numpy arrays]
        filepath = filepath of location to save the object
        """
        time, rfu, voltage, current = data
        with open(filepath, 'w') as file_out:
            file_out.write(header)
            file_out.write('time_data,rfu,voltage,current\n')
            for t, r, v, c in zip(time, rfu, voltage, current):
                file_out.write(f"{t},{r},{v},{c}\n")

    @staticmethod
    def get_data(system):
        """
        Retrieves the data from a CE System
        :param system:
        :return:
        """
        detector_data = system.detector.get_data()
        power_data = system.high_voltage.get_data()
        power_time = power_data['time_data']
        detector_time = detector_data['time_data']
        # We need a power data to be as long as the detector data, so we add in one more time_data pair to be sure
        # thus replace the start and stop of the power time to match the detector time if the detector time has
        # more data points
        if power_time[-1] < detector_time[-1]:
            power_time[-1] = detector_time[-1]
        if power_time[0] > detector_time[0]:
            power_time[0] = detector_time[0]

        # We need to interpolate data values from the power supply (Power supply may not be sampled at the same rate)
        if len(power_time) > 0:
            for channel in ['voltage', 'current']:
                if len(power_data[channel]) > 3:
                    fx = interp1d(power_time, power_data[channel], kind='cubic')
                    power_data[channel] = fx(detector_time)

                else:
                    power_data[channel] = np.pad(power_data[channel],
                                                 (0, len(detector_time) - len(power_data[channel])),
                                                 'edge')

        return [detector_time, detector_data['rfu'], power_data['voltage'], power_data['current']]


class SlideSingleCell:
    """ Create a power point presentation for for single cell data. Each
    slide has the separation, lysis video, peak analysis and fluoresence image
    for every cell.
    """

    def __init__(self, save_to_location='', method=''):
        self.prs = Presentation()
        self.slide_layout = self.prs.slide_layouts[6]
        self.save_to_location = save_to_location

    def add_slide(self, egram_filename, lysis_gif=None, fluor_image=None, separation_info=None, peak_lut=None):

        slide = self.prs.slides.add_slide(self.slide_layout)
        slide_w, slide_h, pad = 10, 7.5, 0.25
        graph_w, graph_h = slide_w / 2, slide_h / 2

        # Add the Peak Integration table
        pks_rfu, pks_time = generate_table(egram_filename, slide, x=pad, y=graph_h + pad, cx=graph_w,
                                           cy=graph_h - 2 * pad)

        # Add the egram chart
        chart_name = self.generate_chart(egram_filename, pks_rfu, pks_time, width=graph_w, height=graph_h)
        slide.shapes.add_picture(chart_name, left=Inches(pad), top=(pad), height=Inches(graph_h))

        # Add the GIF
        if lysis_gif is not None:
            slide.shapes.add_picture(lysis_gif, left=Inches(pad * 2 + graph_w), top=Inches(pad),
                                     height=Inches(graph_h / 2))

        # Add the Fluor Image
        if fluor_image is not None:
            slide.shapes.add_picture(fluor_image, left=Inches(pad * 2 + graph_w), top=Inches(pad + graph_w / 2),
                                     height=Inches(graph_h / 2))

        # Add the separation INfo
        if separation_info is not None:
            tb = slide.shapes.add_textbox(left=Inches(pad * 2 + graph_w), top=Inches(pad * 2 + graph_h),
                                          height=Inches(graph_h / 2), width=Inches(graph_w))
            tf = tb.text_frame
            p = tf.add_paragraph()
            p.text = separation_info
            p.font.size = Pt(8)

    def save_presentation(self):
        self.prs.save(self.save_to_location)

    @staticmethod
    def generate_chart(data_file, peaks_rfu, peaks_time, width=6, height=4, dpi=300):
        fig, ax = plt.subplots()
        ax.set_ylabel('Signal (V)')
        ax.set_xlabel('Time (s)')
        ax.set_title(data_file.split('\\')[-1].strip('.csv'))
        data = pd.read_csv(data_file)

        if 'filt_rfu' in data.columns:
            rfu = data['filt_rfu']
        else:
            rfu = data['rfu']

        _, background = Electropherogram.find_noise(rfu)

        ax.plot(data['time_data'], rfu - background)
        for y, x in zip(peaks_rfu, peaks_time):
            ax.fill_between(x, [0] * len(x), y)

        plt.savefig('temp_chart.png', width=width, height=height, dpi=dpi)

        return 'temp_chart.png'


def generate_table(data_file, slide, x, y, cx, cy):
    """
    Generate table populated with peaks at at least 10x the noise.
    Kinda messy please forgive. Ideally the peak information would be passed in as a
    dataframe or something and not calculated here. But until that is done....
    """

    data = pd.read_csv(data_file)

    if 'filt_rfu' in data.columns:
        rfu = data['filt_rfu']
    else:
        rfu = Electropherogram.filter_butter(data, 1.5, 2)

    # Get the electropherogram Noise and background. Find the peaks.
    noise, background = Electropherogram.find_noise(rfu)
    peaks, properties = find_peaks(rfu, height=10 * noise + background, width=10, rel_height=0.8)
    rfu = rfu - background

    # gather the performance metrics for each peak
    ca = []
    snr = []
    m_time = []
    pks_rfu = []
    pks_time = []
    for pk, start_id, stop_id in zip(peaks, properties['left_ips'], properties['right_ips']):
        start_id, stop_id = [int(x) for x in [start_id, stop_id]]
        pk_rfu, pk_time = rfu[start_id: stop_id], data['time_data'].values[start_id:stop_id]
        pks_rfu.append(pk_rfu)
        pks_time.append(pk_time)
        m_time.append(Electropherogram.get_m1(pk_rfu, pk_time))
        ca.append(Electropherogram.get_corrected_area(pk_rfu, pk_time, m_time[-1]))
        snr.append(max(pk_rfu) / noise)

    # Start Constructing the Table
    headers = ['Migration Time', 'CA', '% Area', 'SNR']
    shape = slide.shapes.add_table(rows=len(peaks) + 1, cols=len(headers),
                                   left=Inches(x), top=Inches(y), width=Inches(cx), height=Inches(cy))
    tbl = shape.table

    # Populate the Headers
    for ix, head in enumerate(headers):
        tbl.cell(0, ix).text = head

    # Populate the Table Information
    row = 1
    total_ca = np.sum(ca)
    for s, m, c in zip(snr, m_time, ca):
        tbl.cell(row, 0).text = f'{m:.1f}'
        tbl.cell(row, 1).text = f'{c:.4e}'
        tbl.cell(row, 2).text = f'{c / total_ca * 100:.2f}'
        tbl.cell(row, 3).text = f'{s:.1f}'
        row += 1
    change_table_font_size(tbl, size=8)

    # Return the Peak information that can be used later on
    return pks_rfu, pks_time


def change_table_font_size(tbl, size=8):
    # Change the font size for each cell
    for cell in iter_cells(tbl):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell